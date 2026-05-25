"""
Gera uma apresentação PowerPoint (.pptx) com o resumo mensal da frota TNORTE.

Usa a biblioteca python-pptx para criar slides programaticamente.
A vantagem: o PPTX gerado é 100% editável no PowerPoint.
"""

import io
from typing import Optional
import pandas as pd

# Tenta importar o python-pptx; se não estiver instalado, marca como indisponível
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_DISPONIVEL = True
except ImportError:
    PPTX_DISPONIVEL = False

# Cores em formato RGB (três bytes: R, G, B) para o python-pptx
if PPTX_DISPONIVEL:
    AZUL    = RGBColor(0x1E, 0x29, 0x5C)   # azul marinho TNORTE
    VERMELHO= RGBColor(0xC8, 0x1D, 0x25)   # vermelho TNORTE
    BRANCO  = RGBColor(0xFF, 0xFF, 0xFF)
    CINZA   = RGBColor(0xF5, 0xF5, 0xF7)
    AZUL2   = RGBColor(0x2D, 0x3F, 0x80)   # azul ligeiramente mais claro para cards
else:
    AZUL = VERMELHO = BRANCO = CINZA = AZUL2 = None


def gerar_pptx(kpis: dict, df_placas: pd.DataFrame, periodo: str) -> Optional[bytes]:
    """
    Gera o PPTX e retorna os bytes do arquivo (para st.download_button).

    kpis      : dicionário retornado por calcular_kpis_gerais()
    df_placas : DataFrame retornado por agregar_por_placa()
    periodo   : string como 'Abril/2026'
    """
    if not PPTX_DISPONIVEL:
        return None

    prs = Presentation()
    # Formato 16:9 (widescreen — padrão atual para apresentações)
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    layout_branco = prs.slide_layouts[6]  # layout sem placeholders

    # Slide 1 — Capa
    _slide_capa(prs.slides.add_slide(layout_branco), periodo, prs)

    # Slide 2 — KPIs Gerais
    _slide_kpis(prs.slides.add_slide(layout_branco), kpis, periodo, prs)

    # Slide 3 — Top Eficiência
    _slide_top_eficiencia(prs.slides.add_slide(layout_branco), df_placas, periodo, prs)

    # Salva em memória (BytesIO) — sem criar arquivo temporário no disco
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


# ─── FUNÇÕES AUXILIARES ───────────────────────────────────────────────────────

def _fundo_azul(slide, prs) -> None:
    """Define o fundo do slide como azul marinho via fill nativo."""
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = AZUL


def _textbox(slide, texto: str, x, y, w, h,
             tamanho: int, negrito: bool = False,
             cor=None, alinhamento=None):
    """
    Adiciona um textbox ao slide com as configurações fornecidas.
    Retorna o shape criado.
    """
    if alinhamento is None:
        alinhamento = PP_ALIGN.LEFT
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = alinhamento
    run = p.add_run()
    run.text = texto
    run.font.size  = Pt(tamanho)
    run.font.bold  = negrito
    if cor:
        run.font.color.rgb = cor
    return box


def _slide_capa(slide, periodo: str, prs) -> None:
    """Slide de capa: logo TNORTE, título e slogan."""
    _fundo_azul(slide, prs)

    _textbox(slide, "TNORTE",
             Inches(1), Inches(1.8), Inches(11), Inches(1.6),
             tamanho=60, negrito=True, cor=BRANCO, alinhamento=PP_ALIGN.CENTER)

    _textbox(slide, f"Gestão de Frota — {periodo}",
             Inches(1), Inches(3.5), Inches(11), Inches(1),
             tamanho=26, negrito=False, cor=VERMELHO, alinhamento=PP_ALIGN.CENTER)

    _textbox(slide, "SINTONIA — TODOS NA MESMA FREQUÊNCIA",
             Inches(1), Inches(5.0), Inches(11), Inches(0.8),
             tamanho=13, negrito=False, cor=RGBColor(0xAA, 0xAA, 0xAA),
             alinhamento=PP_ALIGN.CENTER)


def _slide_kpis(slide, kpis: dict, periodo: str, prs) -> None:
    """Slide com os 4 KPIs principais sobre fundo azul marinho."""
    from core.metrics import fmt_brl, fmt_num

    _fundo_azul(slide, prs)

    # Título do slide
    _textbox(slide, f"Resultados Mensais — {periodo}",
             Inches(0.4), Inches(0.2), Inches(12), Inches(0.7),
             tamanho=22, negrito=True, cor=BRANCO)

    # Dados para os 4 cards
    cards = [
        ("VOLUME ABASTECIDO",   f"{fmt_num(kpis['total_litros'], 2)} L"),
        ("INVESTIMENTO TOTAL",  fmt_brl(kpis['total_valor'])),
        ("DISTÂNCIA PERCORRIDA", f"{fmt_num(kpis['total_km'])} km"),
        ("PREÇO MÉDIO / LITRO", fmt_brl(kpis['preco_medio'])),
    ]

    for i, (label, valor) in enumerate(cards):
        x = Inches(0.3 + i * 3.2)
        y = Inches(1.4)

        # Cria o retângulo de fundo do card usando add_shape (tipo 1 = retângulo)
        card = slide.shapes.add_shape(1, x, y, Inches(3.0), Inches(2.6))
        card.fill.solid()
        card.fill.fore_color.rgb = AZUL2
        card.line.fill.background()  # sem borda

        # Label do KPI
        _textbox(slide, label,
                 x + Inches(0.1), y + Inches(0.2), Inches(2.8), Inches(0.5),
                 tamanho=9, negrito=True,
                 cor=RGBColor(0xBB, 0xBB, 0xBB), alinhamento=PP_ALIGN.CENTER)

        # Valor do KPI
        _textbox(slide, valor,
                 x + Inches(0.05), y + Inches(0.85), Inches(2.9), Inches(1.2),
                 tamanho=17, negrito=True, cor=BRANCO, alinhamento=PP_ALIGN.CENTER)

    # Rodapé
    _textbox(slide, f"TNORTE  |  {kpis['n_abastecimentos']} abastecimentos  |  {kpis['n_placas']} placas",
             Inches(0.4), Inches(6.8), Inches(12), Inches(0.5),
             tamanho=10, cor=RGBColor(0x88, 0x88, 0x88), alinhamento=PP_ALIGN.CENTER)


def _slide_top_eficiencia(slide, df_placas: pd.DataFrame, periodo: str, prs) -> None:
    """Slide com o Top 5 de eficiência geral."""
    from core.metrics import fmt_num

    _fundo_azul(slide, prs)

    _textbox(slide, f"Top 5 Eficiência — {periodo}",
             Inches(0.4), Inches(0.2), Inches(12), Inches(0.7),
             tamanho=22, negrito=True, cor=BRANCO)

    top5 = df_placas[df_placas['MediaKmL'] > 0].nlargest(5, 'MediaKmL')

    medalhas = ["🥇", "🥈", "🥉", "4°", "5°"]

    for i, (_, row) in enumerate(top5.iterrows()):
        y = Inches(1.2 + i * 1.0)
        tipo_txt = str(row.get('Tipo', '')).strip()
        linha = f"{medalhas[i]}  {row['Placa']}  ({tipo_txt})  —  {fmt_num(row['MediaKmL'], 2)} km/L"
        cor_linha = VERMELHO if i == 0 else BRANCO
        _textbox(slide, linha,
                 Inches(1.0), y, Inches(11), Inches(0.7),
                 tamanho=16, negrito=(i == 0), cor=cor_linha)
